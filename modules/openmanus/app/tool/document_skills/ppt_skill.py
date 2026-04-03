from __future__ import annotations

from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .common import (
    ensure_image_path,
    ensure_parent_dir,
    normalize_align,
    parse_bool,
    parse_length_inch,
    parse_length_pt,
    resolve_style,
)
from .css_parser import parse_css_block_from_spec


def _to_pp_align(value: str) -> PP_ALIGN:
    normalized = normalize_align(value, "left")
    if normalized == "center":
        return PP_ALIGN.CENTER
    if normalized == "right":
        return PP_ALIGN.RIGHT
    if normalized == "justify":
        return PP_ALIGN.JUSTIFY
    return PP_ALIGN.LEFT


def _apply_text_style(paragraph: Any, run: Any, style: dict[str, Any]) -> None:
    font = run.font

    if style.get("font_name"):
        font.name = str(style["font_name"])
    if style.get("font_size") is not None:
        font.size = Pt(parse_length_pt(style.get("font_size"), 18.0))
    if style.get("font_color"):
        rgb = str(style["font_color"])
        font.color.rgb = RGBColor.from_string(rgb)

    font.bold = parse_bool(style.get("bold"), False)
    font.italic = parse_bool(style.get("italic"), False)
    font.underline = parse_bool(style.get("underline"), False)

    paragraph.alignment = _to_pp_align(str(style.get("align", "left")))

    if style.get("line_height") is not None:
        paragraph.line_spacing = float(style.get("line_height"))
    if style.get("space_before") is not None:
        paragraph.space_before = Pt(parse_length_pt(style.get("space_before"), 0.0))
    if style.get("space_after") is not None:
        paragraph.space_after = Pt(parse_length_pt(style.get("space_after"), 0.0))


def _add_textbox(slide: Any, item: dict[str, Any], style: dict[str, Any], page_width: float, page_height: float) -> None:
    x = parse_length_inch(item.get("x"), 0.8)
    y = parse_length_inch(item.get("y"), 0.8)

    default_width = max(1.0, page_width - x - 0.8)
    default_height = max(0.5, page_height - y - 0.8)
    width = parse_length_inch(item.get("width"), default_width)
    height = parse_length_inch(item.get("height"), default_height)

    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    text_frame = shape.text_frame
    text_frame.clear()

    text = str(item.get("text", ""))
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text

    if parse_bool(item.get("bullet"), False):
        paragraph.level = int(item.get("bullet_level", 0) or 0)

    _apply_text_style(paragraph, run, style)


def _add_image(slide: Any, item: dict[str, Any], page_width: float, page_height: float) -> None:
    image_path = ensure_image_path(str(item["path"]))

    x = parse_length_inch(item.get("x"), 0.8)
    y = parse_length_inch(item.get("y"), 0.8)

    width_default = max(1.5, page_width - x - 0.8)
    height_default = max(1.0, page_height - y - 0.8)

    width_value = item.get("width")
    height_value = item.get("height")
    width = Inches(parse_length_inch(width_value, width_default)) if width_value is not None else None
    height = Inches(parse_length_inch(height_value, height_default)) if height_value is not None else None

    slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), width=width, height=height)


def _slide_items(slide_spec: dict[str, Any]) -> list[dict[str, Any]]:
    items = slide_spec.get("items")
    if isinstance(items, list):
        return [item for item in items if isinstance(item, dict)]

    elements = slide_spec.get("elements")
    if isinstance(elements, list):
        return [item for item in elements if isinstance(item, dict)]

    inferred: list[dict[str, Any]] = []
    if slide_spec.get("title"):
        inferred.append(
            {
                "type": "title",
                "text": str(slide_spec["title"]),
                "x": 0.8,
                "y": 0.6,
                "width": 11.8,
                "height": 1.0,
                "style_ref": "title",
            }
        )
    if slide_spec.get("body"):
        inferred.append(
            {
                "type": "text",
                "text": str(slide_spec["body"]),
                "x": 0.8,
                "y": 1.8,
                "width": 11.8,
                "height": 4.8,
                "style_ref": "body",
            }
        )
    return inferred


def create_pptx_from_spec(spec: dict[str, Any], output_path: str) -> str:
    presentation = Presentation()

    page = spec.get("page") if isinstance(spec.get("page"), dict) else {}
    page_width_in = parse_length_inch(page.get("width"), 13.333)
    page_height_in = parse_length_inch(page.get("height"), 7.5)

    presentation.slide_width = Inches(page_width_in)
    presentation.slide_height = Inches(page_height_in)

    default_style = spec.get("default_style") if isinstance(spec.get("default_style"), dict) else {}
    css_styles = parse_css_block_from_spec(spec)

    slides = spec.get("slides") if isinstance(spec.get("slides"), list) else []
    if not slides:
        slides = [
            {
                "title": spec.get("title", "Presentation"),
                "body": spec.get("content", ""),
            }
        ]

    requested_layout = int(spec.get("layout_index", 6) or 6)
    layout_index = requested_layout if 0 <= requested_layout < len(presentation.slide_layouts) else 6

    for slide_spec in slides:
        if not isinstance(slide_spec, dict):
            continue

        slide = presentation.slides.add_slide(presentation.slide_layouts[layout_index])

        slide_style = resolve_style(
            default_style=default_style,
            css_styles=css_styles,
            style_ref=slide_spec.get("style_ref"),
            inline_style=slide_spec.get("style") if isinstance(slide_spec.get("style"), dict) else None,
        )
        background = slide_style.get("background_color")
        if background:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor.from_string(str(background))

        for item in _slide_items(slide_spec):
            item_type = str(item.get("type", "text")).lower()
            item_style = resolve_style(
                default_style=slide_style,
                css_styles=css_styles,
                style_ref=item.get("style_ref"),
                inline_style=item.get("style") if isinstance(item.get("style"), dict) else None,
            )

            if item_type in {"text", "textbox", "title", "subtitle", "paragraph"}:
                _add_textbox(slide, item, item_style, page_width_in, page_height_in)
                continue

            if item_type in {"image", "picture"} and item.get("path"):
                _add_image(slide, item, page_width_in, page_height_in)

    destination = ensure_parent_dir(output_path)
    presentation.save(str(destination))
    return str(destination)
